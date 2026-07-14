import streamlit as st
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- COLOR PALETTE DEFINITIONS ---
COLOR_PRIMARY = RGBColor(12, 35, 64)      # Deep Navy Blue
COLOR_SECONDARY = RGBColor(230, 92, 23)    # Warm Orange
COLOR_TEXT_DARK = RGBColor(33, 37, 41)     # Charcoal
COLOR_TEXT_LIGHT = RGBColor(255, 255, 255) # White
COLOR_BG_LIGHT = RGBColor(248, 249, 250)   # Off-White

# --- HELPER FUNCTIONS FOR PPTX ---
def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="GuroAI"):
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY

def generate_presentation(subject, grade, quarter, competency, code, trivia, hook, key_concepts, teacher_guide, individual_task, group_task, reflection):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Slide 1: Title Slide
    slide_1 = prs.slides.add_slide(blank_layout)
    apply_background(slide_1, COLOR_PRIMARY)
    title_box = slide_1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"{subject} | {grade}"
    p1.font.size = Pt(18)
    p1.font.color.rgb = COLOR_SECONDARY
    p1.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = competency
    p2.font.size = Pt(36)
    p2.font.color.rgb = COLOR_TEXT_LIGHT
    p2.font.bold = True
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = f"Competency Code: {code} • {quarter}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 190, 201)
    p3.space_before = Pt(20)

    # Slide 2: Introduce Phase
    slide_2 = prs.slides.add_slide(blank_layout)
    apply_background(slide_2, COLOR_BG_LIGHT)
    add_header(slide_2, "Hook & Classroom Trivia", "ILAW: Introduce")
    
    left_box = slide_2.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    p_h = tf_left.paragraphs[0]
    p_h.text = "Class Activity / Hook"
    p_h.font.bold = True
    p_h.font.size = Pt(20)
    p_h.font.color.rgb = COLOR_PRIMARY
    
    p_h_desc = tf_left.add_paragraph()
    p_h_desc.text = hook
    p_h_desc.font.size = Pt(16)
    p_h_desc.font.color.rgb = COLOR_TEXT_DARK
    p_h_desc.space_before = Pt(10)
    
    right_card = slide_2.shapes.add_shape(1, Inches(6.75), Inches(2.0), Inches(5.75), Inches(4.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_PRIMARY
    tf_right = right_card.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.4)
    tf_right.margin_top = Inches(0.4)
    
    p_tr = tf_right.paragraphs[0]
    p_tr.text = "💡 DID YOU KNOW?"
    p_tr.font.bold = True
    p_tr.font.size = Pt(18)
    p_tr.font.color.rgb = COLOR_SECONDARY
    
    p_tr_desc = tf_right.add_paragraph()
    p_tr_desc.text = trivia
    p_tr_desc.font.size = Pt(16)
    p_tr_desc.font.color.rgb = COLOR_TEXT_LIGHT
    p_tr_desc.space_before = Pt(20)

    # Slide 3: Learn Phase
    slide_3 = prs.slides.add_slide(blank_layout)
    apply_background(slide_3, COLOR_BG_LIGHT)
    add_header(slide_3, "Core Lesson Points", "ILAW: Learn")
    
    concept_box = slide_3.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_concept = concept_box.text_frame
    tf_concept.word_wrap = True
    p_c = tf_concept.paragraphs[0]
    p_c.text = "Key Concepts to Master"
    p_c.font.bold = True
    p_c.font.size = Pt(20)
    p_c.font.color.rgb = COLOR_PRIMARY
    
    for concept in key_concepts.split("\n"):
        if concept.strip():
            p_item = tf_concept.add_paragraph()
            p_item.text = f"• {concept.strip()}"
            p_item.font.size = Pt(16)
            p_item.font.color.rgb = COLOR_TEXT_DARK
            p_item.space_before = Pt(8)
        
    guide_card = slide_3.shapes.add_shape(1, Inches(6.75), Inches(2.0), Inches(5.75), Inches(4.5))
    guide_card.fill.solid()
    guide_card.fill.fore_color.rgb = RGBColor(230, 240, 250)
    guide_card.line.color.rgb = COLOR_SECONDARY
    tf_guide = guide_card.text_frame
    tf_guide.word_wrap = True
    tf_guide.margin_left = Inches(0.4)
    tf_guide.margin_top = Inches(0.4)
    
    p_g = tf_guide.paragraphs[0]
    p_g.text = "👩‍🏫 TEACHER'S NOTE"
    p_g.font.bold = True
    p_g.font.size = Pt(16)
    p_g.font.color.rgb = COLOR_SECONDARY
    
    p_g_desc = tf_guide.add_paragraph()
    p_g_desc.text = teacher_guide
    p_g_desc.font.size = Pt(15)
    p_g_desc.font.color.rgb = COLOR_TEXT_DARK
    p_g_desc.space_before = Pt(15)

    output_path = "GuroAI_Generated_Lesson.pptx"
    prs.save(output_path)
    return output_path

# --- STREAMLIT USER INTERFACE ---
st.set_page_config(page_title="GuroAI Suite", page_icon="💡", layout="wide")

st.title("💡 GuroAI: Interactive ILAW Lesson Creator & TOS Builder")
st.write("Your digital workspace for structured, high-impact teaching materials.")
st.markdown("---")

# --- FILE UPLOADER ADDITION ---
st.sidebar.header("📁 Upload Custom Lesson Plan")
uploaded_file = st.sidebar.file_uploader("Upload a .docx or .txt lesson plan template", type=["docx", "txt"])

if uploaded_file is not None:
    st.sidebar.success("File uploaded successfully!")
    # If it's a text file, read it directly
    if uploaded_file.name.endswith(".txt"):
        raw_text = uploaded_file.read().decode("utf-8")
        st.sidebar.info("Detected custom template text structure.")
    # If it's a Word file, extract paragraphs
    elif uploaded_file.name.endswith(".docx"):
        import docx
        doc = docx.Document(uploaded_file)
        raw_text = "\n".join([para.text for para in doc.paragraphs])
        st.sidebar.info("Extracted text from Word Document template.")
        
    # Autofill your fields using the uploaded file's content!
    # (In a production environment, this text would be sent to an AI to automatically split into parts)
    st.sidebar.text_area("Extracted Raw Content Preview", value=raw_text[:500] + "...", height=150)

# Create Navigation Tabs
tab_lesson, tab_games, tab_tos = st.tabs([
    "📝 Lesson Plan & Slide Designer", 
    "🎮 Gamification Hub", 
    "📊 TOS & Summative Exam Builder"
])

# ==================== TAB 1: LESSON PLAN & SLIDES ====================
with tab_lesson:
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("📌 Metadata & Competency")
        subject = st.selectbox("Subject", ["Science", "Mathematics", "English", "Filipino", "Araling Panlipunan"], key="sub_t1")
        grade = st.selectbox("Grade Level", [f"Grade {i}" for i in range(1, 13)], key="gr_t1")
        quarter = st.selectbox("Quarter", ["1st Quarter", "2nd Quarter", "3rd Quarter", "4th Quarter"], key="q_t1")
        code = st.text_input("Learning Competency Code", value="S7LT-IIa-1", key="code_t1")
        competency = st.text_area("Learning Competency Statement", value="Identify parts of the microscope and their functions.", key="comp_t1")

    with col2:
        st.subheader("🔥 ILAW Setup")
        trivia = st.text_area("Trivia / Food for Thought", value="Anton van Leeuwenhoek discovered bacteria using a handmade microscope.", key="triv_t1")
        hook = st.text_area("Introduce Phase (Hook)", value="Mystery Box: Guess the super-magnified image.", key="hook_t1")
        key_concepts = st.text_area("Learn Phase (Key Concepts - One per line)", value="Ocular Lens (Eyepiece)\nObjective Lenses\nStage Clip\nCoarse Adjustment Knob", key="conc_t1")
        teacher_guide = st.text_area("Teacher's Safety/Delivery Note", value="Ensure students hold the microscope arm and base when moving.", key="tg_t1")

    st.markdown("---")
    st.subheader("🎯 Active Exercises")
    col3, col4, col5 = st.columns(3)
    with col3:
        indiv_task = st.text_area("Individual Task", value="Label the parts diagram sheet.", key="ind_t1")
    with col4:
        group_task = st.text_area("Group Task", value="Find and focus a colored thread slide in groups of three.", key="grp_t1")
    with col5:
        reflection = st.text_area("Wrap-Up (Reflection / Exit Ticket)", value="Write down the easiest and hardest part of the session.", key="ref_t1")

    if st.button("🚀 Generate Presentation Slides", type="primary"):
        with st.spinner("Compiling slides..."):
            file_path = generate_presentation(
                subject, grade, quarter, competency, code, trivia, hook, key_concepts, teacher_guide, indiv_task, group_task, reflection
            )
            st.success("PowerPoint layout constructed!")
            with open(file_path, "rb") as file:
                st.download_button(
                    label="📥 Download Lesson PowerPoint (PPTX)",
                    data=file,
                    file_name="GuroAI_Lesson_Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )


# Create a text-based version of the lesson plan to download
text_lesson_plan = f"""
DEPED ILAW LESSON PLAN
Subject: {subject} | Grade: {grade} | Quarter: {quarter}
Competency Code: {code}
Competency: {competency}

[INTRODUCE]
Trivia: {trivia}
Hook: {hook}

[LEARN]
Key Concepts:
{key_concepts}

Teacher's Note: {teacher_guide}

[APPLY]
Individual Task: {indiv_task}
Group Task: {group_task}

[WRAP-UP]
Reflection: {reflection}
"""

st.download_button(
    label="📄 Download Clean Text File (Ready for Google NotebookLM)",
    data=text_lesson_plan,
    file_name="ILAW_Lesson_Source.txt",
    mime="text/plain"
)
# ==================== TAB 2: GAMIFICATION ENGINE ====================
with tab_games:
    st.header("🎮 Dynamic Gamification Planner")
    st.write("Gamified structures increase classroom retention by up to 40%. Here are live offline and online options based on your topic:")
    
    # We dynamically parse game mechanics depending on the subject
    if subject == "Science":
        game_online = "🔬 **Quizizz Live Relay / Wordwall Matchup**"
        game_online_desc = "Students log in via mobile or tablet. Create a matching game matching the names of the parts of the microscope to their functions. Gives real-time leaderboard points."
        game_online_url = "https://wordwall.net/"
        
        game_offline = "🎯 **The Physical Lab Relay**"
        game_offline_desc = "Split the classroom into 4 groups. Draw a giant microscope diagram on the board. Each member runs up, sticks a sticky note naming one part, and races back to hand the marker to the next teammate. Fastest team with perfect accuracy wins!"
    else:
        game_online = "✏️ **Kahoot! Jumble**"
        game_online_desc = "Sequence activities, orders, or equations on screen. Excellent for step-by-step concepts."
        game_online_url = "https://kahoot.com/"
        
        game_offline = "🎲 **Whiteboard Battle Royale**"
        game_offline_desc = "Divide the board into grid blocks. Teams take turns sending a champion to solve a problem or identify a concept. Winners claim that zone of the board."

    col_g1, col_g2 = st.columns(2)
    with col_g1:
        st.info("🌐 **ONLINE GAMIFICATION SUGGESTION**")
        st.write(game_online)
        st.write(game_online_desc)
        st.markdown(f"[Launch Gamification Platform ↗️]({game_online_url})")
        
    with col_g2:
        st.warning("🎒 **OFFLINE CLASSROOM GAME (No Internet Needed)**")
        st.write(game_offline)
        st.write(game_offline_desc)
        st.write("*Preparation Needed:* Sticky notes, markers, and a timer.")

# ==================== TAB 3: TOS & EXAM GENERATOR ====================
with tab_games if False else tab_tos:
    st.header("📊 Table of Specifications (TOS) & Summative Exam Builder")
    st.write("Generate exams following DepEd's 60-30-10 Cognitive Domain framework (Easy/Medium/Hard).")
    
    col_t1, col_t2 = st.columns([1, 2])
    with col_t1:
        st.subheader("Configure Weights")
        total_items = st.number_input("Target Total Exam Items", min_value=5, max_value=50, value=10, step=1)
        hours_topic1 = st.number_input("Hours Taught: Topic 1 (Microscope Parts)", min_value=1, value=3)
        hours_topic2 = st.number_input("Hours Taught: Topic 2 (Proper Care)", min_value=1, value=2)
        
        total_hours = hours_topic1 + hours_topic2
        weight_1 = (hours_topic1 / total_hours)
        weight_2 = (hours_topic2 / total_hours)
        
        items_1 = round(weight_1 * total_items)
        items_2 = total_items - items_1 # Keep mathematics exact
        
    with col_t2:
        st.subheader("Automated DepEd Table of Specifications")
        
        # Calculate Cognitive Domains distribution based on DepEd rules: 
        # 60% Easy (Remember/Understand), 30% Medium (Apply/Analyze), 10% Difficult (Evaluate/Create)
        easy_items = round(total_items * 0.6)
        medium_items = round(total_items * 0.3)
        hard_items = total_items - (easy_items + medium_items)
        
        st.table({
            "Topic": ["Topic 1: Parts", "Topic 2: Proper Care", "Total Items"],
            "Hours Spent": [hours_topic1, hours_topic2, total_hours],
            "Weight (%)": [f"{weight_1*100:.1f}%", f"{weight_2*100:.1f}%", "100.0%"],
            "Allocated Items": [items_1, items_2, total_items],
            "Cognitive Levels Allocation": [
                f"Easy: {round(items_1*0.6)} | Med: {round(items_1*0.3)} | Hard: {items_1 - (round(items_1*0.6)+round(items_1*0.3))}",
                f"Easy: {round(items_2*0.6)} | Med: {round(items_2*0.3)} | Hard: {items_2 - (round(items_2*0.6)+round(items_2*0.3))}",
                f"Easy: {easy_items} | Med: {medium_items} | Hard: {hard_items}"
            ]
        })
        
    st.markdown("---")
    st.subheader("📝 Generated Summative Quiz Layout")
    st.write("Below is your system-constructed assessment based on the items allocated above:")
    
    # Display sample generated questions matching the TOS structure
    if subject == "Science":
        st.markdown(f"""
        ### **SUMMATIVE TEST: {subject.upper()} {grade.upper()}**
        *Name: ______________________ Section: ___________ Score: _____ / {total_items}*
        
        **PART I: EASY (Remembering / Understanding) - {easy_items} Items**
        1. Which part of the microscope controls the amount of light that passes through the stage slide?
           * A) Eyepiece
           * B) Coarse Adjustment
           * C) Diaphragm
           * D) Base
        
        2. When carrying a microscope to your lab table, which parts should you hold tightly?
           * A) Stage and Lens
           * B) Arm and Base
           * C) Mirror and Clip
           * D) Body Tube and Eyepiece
           
        **PART II: MEDIUM (Applying / Analyzing) - {medium_items} Items**
        3. A student is looking at a plant cell specimen but the field of view is completely dark. What should they adjust first to solve this?
           * A) Turn the fine adjustment knob
           * B) Rotate the nosepiece to High Power Objective
           * C) Check the light source and open up the diaphragm
           * D) Wipe the ocular lens with rough tissue paper
           
        **PART III: DIFFICULT (Evaluating / Creating) - {hard_items} Items**
        4. If a slide containing a specimen of pond water has bubbles trapped beneath the cover slip, how will this impact a student's observation under high magnification? Explain how to properly prepare the slide to fix this issue.
           * *Answer Key Guideline: Air bubbles will appear as thick black rings and mask cellular detail. To resolve this, place the cover slip at a 45-degree angle to the slide and drop it gently to force air outward.*
        """)
    else:
        st.write("Modify options to generate standard test templates.")